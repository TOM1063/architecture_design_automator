import collections.abc
from pptx import Presentation
import openpyxl

import pandas as pd



#-------------------------------------------main-------------------------------------

def main() :

    # keywords = ["データ", "ビジュアライズ", "ビジュアライ","プログラミング"]
    excel_file_name = 'data/result.xlsx'

    sheet_df = pd.read_excel(excel_file_name,sheet_name="input")
    print(sheet_df)
    keywords_df = sheet_df.iloc[0:,0]
    keywords = keywords_df.to_numpy().tolist()
    print(keywords)


    url_sheet_df = pd.read_excel(excel_file_name,sheet_name="url")
    print(url_sheet_df)
    powerpoint_file_name = str(url_sheet_df.iat[0,0])
    print("power_point_file : " + powerpoint_file_name)


    prs = Presentation(powerpoint_file_name)
    db = get_text_data(prs)
    count_result = counter(db,keywords)
    count_result_df = pd.DataFrame(count_result).transpose()
    count_result_df = pd.concat([count_result_df,pd.DataFrame(count_result_df.sum(axis=1),columns=['Total'])],axis=1)
    # print(count_result_df)

    update_sheet_df = pd.concat([sheet_df, count_result_df], axis=1)
    print(update_sheet_df)

    # with pd.ExcelWriter(excel_file_name, mode='a') as writer:
    #     update_sheet_df.to_excel(writer, sheet_name ='result')

    with pd.ExcelWriter(excel_file_name) as writer:
        sheet_df.to_excel(writer, sheet_name ='input',index=False)
        url_sheet_df.to_excel(writer, sheet_name ='url',index=False)
        update_sheet_df.to_excel(writer, sheet_name ='result',index=False)
#-------------------------------------------main-------------------------------------




#-----------------------------------------function-----------------------------------

# get paragraph array by each slide
def get_text_data(presentation):
    slide_num = len(presentation.slides)
    text_db = []
    for i in range(slide_num):
        slide = presentation.slides[i]
        # print(slide)
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                text = paragraph.text
                # print(text)
                if(text):
                    texts.append(text)
        if(True):
            text_db.append(texts)
    return text_db



def counter(database,keywords_list):
    result = []
    for paragraphs in database:
        count_per_key = []
        for key in keywords_list:
            key_count = 0
            for paragraph in paragraphs:
                count = paragraph.count(key)
                key_count += count
            count_per_key.append(key_count)
        result.append(count_per_key)

    #print(result)
    return result
#-----------------------------------------function-----------------------------------


if __name__ == "__main__":
    main()